"""Generate PAARAL ERD as a PowerPoint slide.

Run: python3 docs/gen_erd.py
Output: docs/PAARAL_ERD.pptx
"""


from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ────────────────────────────────────────────────────────
C_STUDENT   = RGBColor(0x1a, 0x4b, 0x8c)   # DepEd blue
C_SCHOOL    = RGBColor(0x16, 0xa3, 0x4a)   # green
C_DEPED     = RGBColor(0xdc, 0x26, 0x26)   # red
C_SHARED    = RGBColor(0x92, 0x40, 0x0e)   # amber-brown
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK      = RGBColor(0x1a, 0x1d, 0x23)
C_LIGHT_BG  = RGBColor(0xf8, 0xf9, 0xfa)
C_BORDER    = RGBColor(0xe2, 0xe4, 0xe9)
C_LINE      = RGBColor(0x94, 0xa3, 0xb8)
C_HEADER_SV = RGBColor(0xdb, 0xea, 0xff)   # student view header bg
C_HEADER_SH = RGBColor(0xdc, 0xfc, 0xe7)   # school view header bg
C_HEADER_DP = RGBColor(0xfe, 0xe2, 0xe2)   # deped view header bg
C_HEADER_CO = RGBColor(0xff, 0xed, 0xd5)   # shared/core header bg

SLIDE_W = Inches(20)
SLIDE_H = Inches(14)

# ── Entity definitions ────────────────────────────────────────────────────
# Each entity: (name, view_tag, [(field_name, type_hint, notes)])
ENTITIES = {
    "SCHOOL": ("SCHOOL", "core", [
        ("id",               "PK string",  "SCH001…SCH050"),
        ("name",             "string",     ""),
        ("type",             "enum",
         "public | private_esc | private_no_esc"),
        ("sector",           "enum",       "sectarian | non_sectarian"),
        ("region",           "string",     ""),
        ("province",         "string",     ""),
        ("municipality",     "string",     ""),
        ("barangay",         "string",     ""),
        ("postal_code",      "string",     ""),
        ("lat / lng",        "decimal",    ""),
        ("tuition",          "int",        "₱ annual"),
        ("esc_subsidy",      "int",        "₱ annual"),
        ("net_cost",         "int",        "tuition − subsidy"),
        ("slots_total",      "int",        ""),
        ("slots_available",  "int",        ""),
        ("distance_km",      "decimal",    ""),
        ("commute_minutes",  "int",        ""),
        ("esc_rating",       "int",        "0–5 stars"),
        ("admission_category", "string",
         "ESC Partner / Public / Selective"),
    ]),
    "LEARNER": ("LEARNER", "student", [
        ("lrn",              "PK string",  "12-digit DepEd LRN"),
        ("given_name",       "string",     ""),
        ("family_name",      "string",     ""),
        ("barangay",         "string",     ""),
        ("municipality",     "string",     ""),
        ("grade6_school_id", "FK string",  "→ SCHOOL"),
    ]),
    "APPLICATION": ("APPLICATION", "student", [
        ("id",               "PK uuid",    ""),
        ("lrn",              "FK string",  "→ LEARNER"),
        ("status",           "enum",       "draft | submitted | cancelled"),
        ("submitted_at",     "timestamp",  ""),
    ]),
    "WISHLIST_ENTRY": ("WISHLIST_ENTRY", "student", [
        ("id",               "PK uuid",    ""),
        ("application_id",   "FK uuid",    "→ APPLICATION"),
        ("school_id",        "FK string",  "→ SCHOOL"),
        ("rank",             "int",        "1 = first choice"),
    ]),
    "ELIGIBILITY_ASSESSMENT": ("ELIGIBILITY_ASSESSMENT", "student", [
        ("id",               "PK uuid",    ""),
        ("application_id",   "FK uuid",    "→ APPLICATION"),
        ("esc_intent",       "enum",       "yes | no"),
        ("school_type",      "enum",       "public | private | als"),
        ("segs",             "json",
         "4ps, gidca, ip, pwd, special, cbms"),
        ("income",           "enum",
         "poor | low | lower_middle | middle | above"),
        ("employment",       "enum",
         "local | abroad | business | unemployed"),
        ("category",         "enum",       "A | B | C | D | none"),
    ]),
    "SURVEY_RESPONSE": ("SURVEY_RESPONSE", "student", [
        ("id",               "PK uuid",    ""),
        ("application_id",   "FK uuid",    "→ APPLICATION"),
        ("ease",             "int",        "1–5 scale"),
        ("helpful",          "enum",       "yes | somewhat | no"),
        ("concern",          "enum",
         "cost | distance | quality | slots"),
        ("suggestions",      "text",       "optional"),
    ]),
    "DOCUMENT_UPLOAD": ("DOCUMENT_UPLOAD", "student", [
        ("id",               "PK uuid",    ""),
        ("application_id",   "FK uuid",    "→ APPLICATION"),
        ("doc_type",         "string",     ""),
        ("storage_path",     "string",     ""),
        ("uploaded_at",      "timestamp",  ""),
    ]),
    "SCHOOL_USER": ("SCHOOL_USER", "school", [
        ("id",               "PK uuid",    ""),
        ("school_id",        "FK string",  "→ SCHOOL"),
        ("name",             "string",     ""),
        ("role",             "enum",       "staff | esc_committee_member"),
    ]),
    "APPLICATION_REVIEW": ("APPLICATION_REVIEW", "school", [
        ("id",               "PK uuid",    ""),
        ("application_id",   "FK uuid",    "→ APPLICATION"),
        ("school_id",        "FK string",  "→ SCHOOL"),
        ("reviewed_by",      "FK uuid",    "→ SCHOOL_USER"),
        ("decision",         "enum",
         "pending | approved | waitlisted | rejected"),
        ("remarks",          "text",       ""),
        ("reviewed_at",      "timestamp",  ""),
    ]),
    "SLOT_ALLOCATION": ("SLOT_ALLOCATION", "school", [
        ("id",               "PK uuid",    ""),
        ("school_id",        "FK string",  "→ SCHOOL"),
        ("school_year",      "string",     "e.g. 2026-2027"),
        ("slots_allocated",  "int",        ""),
        ("slots_consumed",   "int",        ""),
        ("status",           "enum",       "open | closed | full"),
    ]),
    "REGION": ("REGION", "deped", [
        ("code",             "PK string",  ""),
        ("name",             "string",     ""),
    ]),
    "DIVISION": ("DIVISION", "deped", [
        ("code",             "PK string",  ""),
        ("name",             "string",     ""),
        ("region_code",      "FK string",  "→ REGION"),
    ]),
    "PLANNING_SNAPSHOT": ("PLANNING_SNAPSHOT", "deped", [
        ("id",               "PK uuid",    ""),
        ("region_code",      "FK string",  "→ REGION"),
        ("school_year",      "string",     ""),
        ("projected_g7_enroll", "int",     ""),
        ("public_jhs_capacity", "int",     ""),
        ("esc_slots_needed", "int",        ""),
        ("esc_slots_available", "int",     ""),
        ("generated_at",     "timestamp",  ""),
    ]),
}

# ── Relationships ─────────────────────────────────────────────────────────
# (from_entity, to_entity, label, cardinality_from, cardinality_to)
# cardinality: "1" | "0..1" | "1..*" | "0..*"
RELATIONSHIPS = [
    ("LEARNER", "APPLICATION", "submits", "1", "0..*"),
    ("LEARNER", "SCHOOL", "attended (Gr 6)", "0..*", "0..1"),
    ("APPLICATION", "WISHLIST_ENTRY", "contains", "1", "1..*"),
    ("APPLICATION", "ELIGIBILITY_ASSESSMENT", "has", "1", "0..1"),
    ("APPLICATION", "SURVEY_RESPONSE", "has", "1", "0..1"),
    ("APPLICATION", "DOCUMENT_UPLOAD", "attaches", "1", "0..*"),
    ("APPLICATION", "APPLICATION_REVIEW", "reviewed via", "1", "0..*"),
    ("SCHOOL", "WISHLIST_ENTRY", "appears in", "1", "0..*"),
    ("SCHOOL", "SLOT_ALLOCATION", "has", "1", "0..*"),
    ("SCHOOL", "APPLICATION_REVIEW", "receives", "1", "0..*"),
    ("SCHOOL", "DIVISION", "governed by", "0..*", "1"),
    ("SCHOOL_USER", "SCHOOL", "belongs to", "0..*", "1"),
    ("SCHOOL_USER", "APPLICATION_REVIEW", "makes", "1", "0..*"),
    ("REGION", "DIVISION", "contains", "1", "1..*"),
    ("REGION", "PLANNING_SNAPSHOT", "analyzed in", "1", "0..*"),
]

# ── Layout positions (x, y) in inches ────────────────────────────────────
# Designed for a 20×14 inch slide
POSITIONS = {
    # Core (centre top)
    "SCHOOL":                  (8.2,  0.9),
    # Student view (left column)
    "LEARNER":                 (0.3,  0.9),
    "APPLICATION":             (0.3,  4.2),
    "WISHLIST_ENTRY":          (0.3,  7.6),
    "ELIGIBILITY_ASSESSMENT":  (4.2,  0.9),
    "SURVEY_RESPONSE":         (4.2,  5.2),
    "DOCUMENT_UPLOAD":         (4.2,  8.8),
    # School view (right column)
    "SCHOOL_USER":             (14.2, 0.9),
    "APPLICATION_REVIEW":      (14.2, 4.0),
    "SLOT_ALLOCATION":         (14.2, 8.0),
    # DepEd view (far right)
    "REGION":                  (17.5, 0.9),
    "DIVISION":                (17.5, 3.8),
    "PLANNING_SNAPSHOT":       (17.5, 6.8),
}

VIEW_COLORS = {
    "student": (C_STUDENT,   C_HEADER_SV),
    "school":  (C_SCHOOL,    C_HEADER_SH),
    "deped":   (C_DEPED,     C_HEADER_DP),
    "core":    (C_SHARED,    C_HEADER_CO),
}

ENTITY_W = Inches(2.95)


def field_rows(fields):
    """Return the number of rows in an entity field list.

    Parameters
    ----------
    fields : list
        List of (name, type_hint, notes) tuples.

    Returns
    -------
    int
        Number of fields.
    """
    return len(fields)


def entity_height(fields):
    """Compute the total height of an entity box.

    Parameters
    ----------
    fields : list
        List of (name, type_hint, notes) tuples.

    Returns
    -------
    pptx.util.Length
        Total box height (header + rows + bottom padding).
    """
    header_h = Inches(0.38)
    row_h = Inches(0.22)
    return header_h + row_h * len(fields) + Inches(0.10)


# ── Helpers ───────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_w=Pt(1)):
    """Add a filled rectangle shape to a slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    x, y : pptx.util.Length
        Top-left corner position.
    w, h : pptx.util.Length
        Width and height.
    fill_rgb : pptx.dml.color.RGBColor
        Fill colour.
    line_rgb : pptx.dml.color.RGBColor, optional
        Border colour. No border if None.
    line_w : pptx.util.Length
        Border width.

    Returns
    -------
    pptx.shapes.autoshape.Shape
        The created rectangle shape.
    """
    shape = slide.shapes.add_shape(1, x, y, w, h)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h, font_size=Pt(8),
                 bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
                 wrap=True, italic=False):
    """Add a formatted text box to a slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    text : str
        Text content.
    x, y : pptx.util.Length
        Position of the text box.
    w, h : pptx.util.Length
        Width and height.
    font_size : pptx.util.Length
        Font size.
    bold : bool
        Bold formatting.
    color : pptx.dml.color.RGBColor
        Font colour.
    align : PP_ALIGN
        Paragraph alignment.
    wrap : bool
        Whether to wrap text.
    italic : bool
        Italic formatting.

    Returns
    -------
    pptx.shapes.picture.Picture
        The created text box shape.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2,
             color=C_LINE, width=Pt(1.2), dash=None):
    """Add a connector line between two points on a slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    x1, y1 : pptx.util.Length
        Start point.
    x2, y2 : pptx.util.Length
        End point.
    color : pptx.dml.color.RGBColor
        Line colour.
    width : pptx.util.Length
        Line width.
    dash : pptx.enum.dml.MSO_LINE_DASH_STYLE, optional
        Dash style. Solid if None.

    Returns
    -------
    pptx.shapes.connector.Connector
        The created connector shape.
    """
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    if dash:
        connector.line.dash_style = dash
    return connector


def inches_to_emu(val_in_inches):
    """Convert a value in inches to English Metric Units.

    Parameters
    ----------
    val_in_inches : float
        Value in inches.

    Returns
    -------
    int
        Equivalent value in EMU (914400 EMU per inch).
    """
    return int(val_in_inches * 914400)


# ── Entity box renderer ───────────────────────────────────────────────────

def draw_entity(slide, key, x_in, y_in):
    """Draw a single entity box onto the slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    key : str
        Entity name key in ENTITIES.
    x_in, y_in : float
        Top-left position in inches.

    Returns
    -------
    tuple
        (x, y, w, h) as pptx Length values.
    """
    _, view_tag, fields = ENTITIES[key]
    accent_color, header_bg = VIEW_COLORS[view_tag]

    x = Inches(x_in)
    y = Inches(y_in)
    w = ENTITY_W
    h = entity_height(fields)

    header_h = Inches(0.38)
    row_h = Inches(0.22)

    # Outer border
    add_rect(slide, x, y, w, h, C_WHITE, accent_color, Pt(1.5))

    # Header band
    add_rect(slide, x, y, w, header_h, header_bg)
    # Header accent left strip
    add_rect(slide, x, y, Inches(0.07), header_h, accent_color)

    # Entity name in header
    add_text_box(slide, key, x + Inches(0.14), y + Inches(0.07),
                 w - Inches(0.18), header_h,
                 font_size=Pt(7.5), bold=True, color=accent_color)

    # Divider line under header
    add_rect(slide, x, y + header_h, w, Pt(1), accent_color)

    # Fields
    for i, (fname, ftype, fnote) in enumerate(fields):
        fy = y + header_h + Inches(0.04) + i * row_h
        # Alternating row tint
        if i % 2 == 0:
            add_rect(slide, x + Inches(0.01), fy,
                     w - Inches(0.02), row_h,
                     RGBColor(0xf8, 0xf9, 0xfb))

        # PK/FK tag
        tag_text = ""
        display_fname = fname
        if "PK" in ftype:
            tag_text = "PK"
        elif "FK" in ftype:
            tag_text = "FK"

        if tag_text:
            tag_bg = (
                RGBColor(0xff, 0xf0, 0xd0)
                if tag_text == "FK"
                else RGBColor(0xe0, 0xee, 0xff)
            )
            tag_col = (
                RGBColor(0x92, 0x40, 0x0e)
                if tag_text == "FK"
                else RGBColor(0x1a, 0x4b, 0x8c)
            )
            add_rect(slide, x + Inches(0.06), fy + Inches(0.03),
                     Inches(0.22), Inches(0.16), tag_bg)
            add_text_box(slide, tag_text,
                         x + Inches(0.06), fy + Inches(0.02),
                         Inches(0.22), Inches(0.18),
                         font_size=Pt(5.5), bold=True, color=tag_col,
                         align=PP_ALIGN.CENTER)
            field_x = x + Inches(0.32)
        else:
            field_x = x + Inches(0.10)

        # Field name
        clean_type = ftype.replace("PK ", "").replace("FK ", "")
        add_text_box(slide, display_fname,
                     field_x, fy + Inches(0.02),
                     Inches(1.28), row_h,
                     font_size=Pt(6.8), bold=False, color=C_DARK)

        # Type
        add_text_box(slide, clean_type,
                     x + Inches(1.42), fy + Inches(0.02),
                     Inches(0.82), row_h,
                     font_size=Pt(6), bold=False,
                     color=RGBColor(0x6b, 0x72, 0x80), italic=True)

        # Note
        if fnote:
            add_text_box(slide, fnote,
                         x + Inches(2.26), fy + Inches(0.02),
                         Inches(0.66), row_h,
                         font_size=Pt(5.5), bold=False,
                         color=RGBColor(0x9c, 0xa3, 0xaf), italic=True)

    return x, y, w, h


# ── Relationship line renderer ────────────────────────────────────────────

def entity_centre(key):
    """Return the centre point of an entity box.

    Parameters
    ----------
    key : str
        Entity name key in ENTITIES and POSITIONS.

    Returns
    -------
    tuple
        (cx, cy) as pptx Length values.
    """
    x_in, y_in = POSITIONS[key]
    _, _, fields = ENTITIES[key]
    h_in = entity_height(fields).inches
    cx = Inches(x_in) + ENTITY_W / 2
    cy = Inches(y_in) + Inches(h_in / 2)
    return cx, cy


def entity_bbox(key):
    """Return the bounding box of an entity box.

    Parameters
    ----------
    key : str
        Entity name key in ENTITIES and POSITIONS.

    Returns
    -------
    tuple
        (x, y, w, h) as pptx Length values.
    """
    x_in, y_in = POSITIONS[key]
    _, _, fields = ENTITIES[key]
    h = entity_height(fields)
    return Inches(x_in), Inches(y_in), ENTITY_W, h


def edge_point(key, towards_key):
    """Return border anchor on key facing towards_key.

    Parameters
    ----------
    key : str
        Source entity name.
    towards_key : str
        Target entity name (determines which border to use).

    Returns
    -------
    tuple
        (x, y) anchor point as pptx Length values.
    """
    ax, ay, aw, ah = entity_bbox(key)
    bx, by, bw, bh = entity_bbox(towards_key)
    acx = ax + aw / 2
    acy = ay + ah / 2
    bcx = bx + bw / 2
    bcy = by + bh / 2

    dx = bcx - acx
    dy = bcy - acy

    # Choose exit side based on dominant direction
    if abs(dx) >= abs(dy):
        if dx > 0:
            return ax + aw, acy   # right edge
        else:
            return ax, acy        # left edge
    else:
        if dy > 0:
            return acx, ay + ah   # bottom edge
        else:
            return acx, ay        # top edge


def draw_relationship(slide, from_key, to_key, label):
    """Draw a labelled relationship line between two entities.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    from_key : str
        Source entity name.
    to_key : str
        Target entity name.
    label : str
        Relationship label shown at line midpoint.
    """
    x1, y1 = edge_point(from_key, to_key)
    x2, y2 = edge_point(to_key, from_key)
    add_line(slide, x1, y1, x2, y2, color=C_LINE, width=Pt(1.4))

    # Label at midpoint
    mx = (x1 + x2) / 2
    my = (y1 + y2) / 2
    add_text_box(slide, label,
                 mx - Inches(0.55), my - Inches(0.14),
                 Inches(1.1), Inches(0.22),
                 font_size=Pt(5.8), bold=False,
                 color=RGBColor(0x4b, 0x55, 0x63),
                 align=PP_ALIGN.CENTER)


# ── Legend ────────────────────────────────────────────────────────────────

def draw_legend(slide):
    """Draw the colour-coded view legend at the bottom of the slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    """
    lx = Inches(0.3)
    ly = Inches(12.4)
    lw = Inches(19.4)
    lh = Inches(1.3)

    add_rect(slide, lx, ly, lw, lh, C_LIGHT_BG, C_BORDER, Pt(1))

    items = [
        ("Student View",  C_STUDENT,  C_HEADER_SV),
        ("School View",   C_SCHOOL,   C_HEADER_SH),
        ("DepEd View",    C_DEPED,    C_HEADER_DP),
        ("Shared / Core", C_SHARED,   C_HEADER_CO),
    ]

    add_text_box(slide, "LEGEND", lx + Inches(0.2), ly + Inches(0.1),
                 Inches(1), Inches(0.3),
                 font_size=Pt(7), bold=True, color=C_DARK)

    for i, (label, accent, bg) in enumerate(items):
        ix = lx + Inches(0.2) + i * Inches(2.6)
        iy = ly + Inches(0.42)
        add_rect(slide, ix, iy,
                 Inches(2.3), Inches(0.28), bg, accent, Pt(1.5))
        add_rect(slide, ix, iy, Inches(0.07), Inches(0.28), accent)
        add_text_box(slide, label, ix + Inches(0.14), iy + Inches(0.05),
                     Inches(2.1), Inches(0.22),
                     font_size=Pt(7.5), bold=True, color=accent)

    tag_items = [
        ("PK", "Primary Key",
         RGBColor(0xe0, 0xee, 0xff), RGBColor(0x1a, 0x4b, 0x8c)),
        ("FK", "Foreign Key",
         RGBColor(0xff, 0xf0, 0xd0), RGBColor(0x92, 0x40, 0x0e)),
    ]
    for i, (tag, desc, bg, col) in enumerate(tag_items):
        ix = lx + Inches(10.8) + i * Inches(2.0)
        iy = ly + Inches(0.42)
        add_rect(slide, ix, iy, Inches(0.28), Inches(0.28), bg)
        add_text_box(slide, tag, ix, iy + Inches(0.04),
                     Inches(0.28), Inches(0.22),
                     font_size=Pt(6.5), bold=True, color=col,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, f"= {desc}",
                     ix + Inches(0.32), iy + Inches(0.05),
                     Inches(1.6), Inches(0.22),
                     font_size=Pt(7), bold=False, color=C_DARK)

    # Note about mockup phase
    add_text_box(slide,
                 "⚠  Fields marked with an asterisk (*) are mocked "
                 "in the current pilot phase. Production implementation "
                 "will use real LIS / BEIS data sources.",
                 lx + Inches(0.2), ly + Inches(0.82),
                 lw - Inches(0.4), Inches(0.36),
                 font_size=Pt(6.5), bold=False,
                 color=RGBColor(0x6b, 0x72, 0x80), italic=True)


# ── Section labels ────────────────────────────────────────────────────────

def draw_section_labels(slide):
    """Draw the coloured view-section header bands across the top.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    """
    sections = [
        ("STUDENT VIEW", Inches(0.3),  Inches(0.3), C_STUDENT,  C_HEADER_SV),
        ("CORE",         Inches(7.5),  Inches(0.3), C_SHARED,   C_HEADER_CO),
        ("SCHOOL VIEW",  Inches(13.6), Inches(0.3), C_SCHOOL,   C_HEADER_SH),
        ("DEPED VIEW",   Inches(17.0), Inches(0.3), C_DEPED,    C_HEADER_DP),
    ]
    widths = [Inches(7.0), Inches(5.8), Inches(3.2), Inches(2.6)]

    for (label, lx, ly, col, bg), lw in zip(sections, widths):
        add_rect(slide, lx, ly, lw, Inches(0.32), bg, col, Pt(1.5))
        add_rect(slide, lx, ly, Inches(0.07), Inches(0.32), col)
        add_text_box(slide, label, lx + Inches(0.14), ly + Inches(0.06),
                     lw, Inches(0.24),
                     font_size=Pt(8), bold=True, color=col)


# ── Title ─────────────────────────────────────────────────────────────────

def draw_title(slide):
    """Draw the diagram title at the top of the slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    """
    add_text_box(slide,
                 "PAARAL — Entity-Relationship Diagram"
                 "  |  Student · School · DepEd Views",
                 Inches(0.3), Inches(0.0),
                 Inches(19.4), Inches(0.32),
                 font_size=Pt(10), bold=True, color=C_DARK)


# ── Main ──────────────────────────────────────────────────────────────────

def main():
    """Build the ERD PowerPoint and save it to docs/PAARAL_ERD.pptx."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_WHITE

    draw_title(slide)
    draw_section_labels(slide)

    # Draw all entities
    for key, (x_in, y_in) in POSITIONS.items():
        draw_entity(slide, key, x_in, y_in)

    # pptx z-order: last added = on top.
    # Draw rels first, then entities, on a second fresh slide.
    prs2 = Presentation()
    prs2.slide_width  = SLIDE_W
    prs2.slide_height = SLIDE_H
    blank_layout2 = prs2.slide_layouts[6]
    slide2 = prs2.slides.add_slide(blank_layout2)
    bg2 = slide2.background
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = C_WHITE

    draw_title(slide2)
    draw_section_labels(slide2)

    # Draw relationships first (so they appear under entity boxes)
    for (from_e, to_e, label, _, _) in RELATIONSHIPS:
        draw_relationship(slide2, from_e, to_e, label)

    # Then draw entities on top
    for key, (x_in, y_in) in POSITIONS.items():
        draw_entity(slide2, key, x_in, y_in)

    draw_legend(slide2)

    out_path = (
        "/Users/paulamartinez/paaral-student-view/docs/PAARAL_ERD.pptx"
    )
    prs2.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
